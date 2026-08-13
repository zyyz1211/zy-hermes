#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create a PowerPoint presentation matching the 集团公司 template style.

Template: 集团公司模板.ppt (binary .ppt format)
This script builds a new .pptx with matching fonts, colors, and layout.

Usage:
    python create_ppt.py --output "汇报材料.pptx" --title "主标题" --slides slides.json

Slides JSON format:
[
  {
    "layout": "title" | "section" | "content" | "blank",
    "title": "幻灯片标题",
    "subtitle": "副标题（仅 title 布局）",
    "body": ["要点1", "要点2", ...],
    "note": "备注文字（可选）"
  }
]

Or pass a single markdown string with --markdown:
    # 主标题
    ## 副标题
    ---
    ## 章节标题
    内容段落
    - 列表项
    ---
    ...
"""

import json
import argparse
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ── Template constants ─────────────────────────────────────────────
SLIDE_WIDTH = Emu(12192000)   # 13.33 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.50 inches

# Theme colors (from the actual template)
# Cover/end/section use full #002060 dark navy; content title bars use #2A519E (measured from real output)
COLOR_COVER_BG = RGBColor(0x00, 0x20, 0x60)           # 002060 dark navy (cover, end, section)
COLOR_TITLE_BAR_BG = RGBColor(0x2A, 0x51, 0x9E)       # 2A519E (content slide title bars, per design spec)
COLOR_TITLE_TEXT = RGBColor(0xFF, 0xFF, 0xFF)          # White on dark bg
COLOR_SECTION_TITLE = RGBColor(0x00, 0x70, 0xC0)       # 0070C0 blue
COLOR_SUB_TITLE = RGBColor(0x00, 0x20, 0x60)           # 002060 dark navy
COLOR_BODY_TEXT = RGBColor(0x33, 0x33, 0x33)           # Dark gray
COLOR_ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)         # 0070C0
COLOR_ACCENT_LIGHT_BLUE = RGBColor(0x00, 0xB0, 0xF0)   # 00B0F0
COLOR_ACCENT_RED = RGBColor(0xFF, 0x00, 0x00)          # FF0000
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_TITLE = "微软雅黑"
FONT_BODY = "微软雅黑"
FONT_EN = "Franklin Gothic Medium"

# Title bar height
TITLE_BAR_HEIGHT = Inches(1.0)

# ── Helper functions ───────────────────────────────────────────────

def set_slide_size(prs):
    """Set slide dimensions to match template."""
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT


def _set_font(p, font_name, font_size=None, bold=False, color=None,
              ea_font_name=None):
    """Set both Latin and East Asian font on a paragraph.

    python-pptx's ``p.font.name`` only sets the ``latin`` typeface. For
    Chinese characters (微软雅黑 etc.) we must also write an ``a:ea``
    element on the run properties or the text renders as 宋体.
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    p.font.name = font_name
    if font_size is not None:
        p.font.size = font_size
    if bold:
        p.font.bold = bold
    if color is not None:
        p.font.color.rgb = color
    # Set East Asian typeface via the paragraph's pPr/rPr element
    ea = ea_font_name or font_name
    pPr = p._p.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p._p, qn('a:pPr'))
    rPr = pPr.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.SubElement(pPr, qn('a:rPr'))
    # Remove existing a:ea if any
    for child in list(rPr):
        if child.tag == qn('a:ea'):
            rPr.remove(child)
    ea_elem = etree.SubElement(rPr, qn('a:ea'))
    ea_elem.set('typeface', ea)


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=Pt(14), bold=False, color=COLOR_BODY_TEXT,
                alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    _set_font(p, font_name, font_size=font_size, bold=bold, color=color)
    p.alignment = alignment
    return txBox


def add_title_bar(slide, title_text):
    """Add a dark blue title bar at the top with white text."""
    # Background rectangle
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_WIDTH, TITLE_BAR_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_TITLE_BAR_BG
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.1), SLIDE_WIDTH - Inches(1.0), Inches(0.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    _set_font(p, FONT_TITLE, font_size=Pt(24), bold=True, color=COLOR_TITLE_TEXT)
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_subtitle_bar(slide, subtitle_text, top=Inches(1.2)):
    """Add a section subtitle with dark blue text."""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), top, SLIDE_WIDTH - Inches(1.0), Inches(0.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle_text
    _set_font(p, FONT_TITLE, font_size=Pt(18), bold=True, color=COLOR_SUB_TITLE)
    return txBox


def add_body_text(slide, items, top=Inches(1.9), left=Inches(0.8)):
    """Add body text with bullet points."""
    txBox = slide.shapes.add_textbox(
        left, top, SLIDE_WIDTH - Inches(1.6), SLIDE_HEIGHT - top - Inches(0.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        _set_font(p, FONT_BODY, font_size=Pt(14), color=COLOR_BODY_TEXT)
        p.space_after = Pt(6)
        p.level = 0

        # Support sub-items starting with "  " (2 spaces) or "\t"
        if item.startswith("  ") or item.startswith("\t"):
            p.level = 1
            _set_font(p, FONT_BODY, font_size=Pt(12), color=COLOR_BODY_TEXT)
    return txBox


def add_decorative_line(slide, top=Inches(1.05)):
    """Add a thin decorative line under the title bar."""
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), top, SLIDE_WIDTH - Inches(1.0), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT_BLUE
    line.line.fill.background()
    return line


def add_footer(slide, page_number=1):
    """Add '某工程集团' brand footer and page number."""
    from pptx.util import Pt
    # Left: company name
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.9), Inches(5.0), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "某工程集团"
    _set_font(p, "微软雅黑", font_size=Pt(9), color=RGBColor(0x66, 0x66, 0x66))

    # Right: page number
    txBox2 = slide.shapes.add_textbox(
        Inches(11.0), Inches(6.9), Inches(1.0), Inches(0.3)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_number)
    _set_font(p2, "Arial", font_size=Pt(9), color=RGBColor(0x66, 0x66, 0x66), ea_font_name="Arial")
    p2.alignment = PP_ALIGN.RIGHT


# ── Slide builders ─────────────────────────────────────────────────

def make_title_slide(prs, title, subtitle=""):
    """Create a cover/title slide matching the template style."""
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide = prs.slides.add_slide(slide_layout)

    # Remove existing placeholders
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Background: dark navy
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_COVER_BG

    # Main title
    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), SLIDE_WIDTH - Inches(2.0), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    _set_font(p, FONT_EN, font_size=Pt(36), bold=True, color=COLOR_WHITE, ea_font_name="微软雅黑")
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            Inches(1.0), Inches(3.8), SLIDE_WIDTH - Inches(2.0), Inches(0.8)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        _set_font(p2, FONT_BODY, font_size=Pt(16), color=COLOR_WHITE)
        p2.alignment = PP_ALIGN.CENTER

    # Department name at bottom
    txBox3 = slide.shapes.add_textbox(
        Inches(1.0), Inches(5.5), SLIDE_WIDTH - Inches(2.0), Inches(0.5)
    )
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "经营管理部"
    _set_font(p3, FONT_EN, font_size=Pt(32), bold=True, color=COLOR_WHITE, ea_font_name="微软雅黑")
    p3.alignment = PP_ALIGN.CENTER

    return slide


def make_content_slide(prs, title, subtitle="", body=None, page_number=1):
    """Create a content slide with title bar and body text."""
    slide_layout = prs.slide_layouts[6]  # 空白 layout
    slide = prs.slides.add_slide(slide_layout)

    # Remove placeholders
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Title bar
    add_title_bar(slide, title)

    # Decorative line
    add_decorative_line(slide)

    # Subtitle
    if subtitle:
        add_subtitle_bar(slide, subtitle)

    # Body
    if body:
        body_top = Inches(1.9) if subtitle else Inches(1.3)
        add_body_text(slide, body, top=body_top)

    # Footer
    add_footer(slide, page_number=page_number)

    return slide


def make_section_slide(prs, section_title, items=None):
    """Create a section divider / TOC-style slide."""
    slide_layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(slide_layout)

    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_COVER_BG

    # Section number circle placeholder
    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.5), Inches(1.0), Inches(1.0)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u25cf"
    _set_font(p, FONT_TITLE, font_size=Pt(36), color=COLOR_WHITE)
    p.alignment = PP_ALIGN.CENTER

    # Section title
    txBox2 = slide.shapes.add_textbox(
        Inches(2.0), Inches(1.5), SLIDE_WIDTH - Inches(3.0), Inches(1.0)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = section_title
    _set_font(p2, FONT_EN, font_size=Pt(36), bold=True, color=COLOR_WHITE, ea_font_name="微软雅黑")
    p2.alignment = PP_ALIGN.LEFT

    # Items
    if items:
        txBox3 = slide.shapes.add_textbox(
            Inches(2.0), Inches(3.0), SLIDE_WIDTH - Inches(3.0), Inches(3.0)
        )
        tf3 = txBox3.text_frame
        for i, item in enumerate(items):
            if i == 0:
                pp = tf3.paragraphs[0]
            else:
                pp = tf3.add_paragraph()
            pp.text = f"  {item}"
            _set_font(pp, FONT_BODY, font_size=Pt(20), color=COLOR_WHITE)
            pp.space_after = Pt(8)

    return slide


def make_end_slide(prs):
    """Create an end/thank-you slide."""
    slide_layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(slide_layout)

    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_COVER_BG

    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), SLIDE_WIDTH - Inches(2.0), Inches(2.0)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "汇报结束，谢谢！"
    _set_font(p, FONT_EN, font_size=Pt(48), bold=True, color=COLOR_WHITE, ea_font_name="微软雅黑")
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Thanks For Your Attention"
    _set_font(p2, FONT_EN, font_size=Pt(48), bold=True, color=COLOR_WHITE, ea_font_name="微软雅黑")
    p2.alignment = PP_ALIGN.CENTER

    return slide


# ── Markdown parser ────────────────────────────────────────────────

def parse_markdown(md_text):
    """Parse a simple markdown into slide definitions.

    Format:
        # Slide Title
        ## Subtitle
        ---
        ## Section Title
        - Item 1
        - Item 2
        ---
    """
    slides = []
    current = {"layout": "content", "title": "", "body": []}
    in_content = False

    for line in md_text.strip().split("\n"):
        line_stripped = line.strip()

        # Slide separator
        if line_stripped == "---":
            if current["title"] or current["body"]:
                slides.append(current)
            current = {"layout": "content", "title": "", "body": []}
            in_content = False
            continue

        # Title (H1)
        if line_stripped.startswith("# ") and not line_stripped.startswith("## "):
            current["title"] = line_stripped[2:].strip()
            in_content = False
            continue

        # Subtitle (H2) — only on first slide
        if line_stripped.startswith("## "):
            if not current["title"]:
                current["title"] = line_stripped[3:].strip()
            elif not current.get("subtitle"):
                current["subtitle"] = line_stripped[3:].strip()
            else:
                current["body"].append(line_stripped[3:].strip())
            in_content = True
            continue

        # Bullet items
        if line_stripped.startswith("- "):
            current["body"].append(line_stripped[2:].strip())
            in_content = True
            continue

        # Plain text (body content)
        if line_stripped and in_content:
            current["body"].append(line_stripped)

    # Last slide
    if current["title"] or current["body"]:
        slides.append(current)

    return slides


# ── Main ───────────────────────────────────────────────────────────

def create_presentation(slides, output_path, template_path=None):
    """Create a PPTX from slide definitions."""
    if template_path and Path(template_path).exists() and template_path.lower().endswith('.pptx'):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # Remove default slide
        if prs.slides:
            rId = prs.slides._sldIdLst[0].get(qn('r:id'))
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    set_slide_size(prs)

    for i, slide_def in enumerate(slides):
        layout = slide_def.get("layout", "content")
        title = slide_def.get("title", "")
        subtitle = slide_def.get("subtitle", "")
        body = slide_def.get("body", [])

        if layout == "title" and i == 0:
            make_title_slide(prs, title, subtitle)
        elif layout == "section":
            make_section_slide(prs, title, body if body else None)
        elif layout == "end":
            make_end_slide(prs)
        else:
            make_content_slide(prs, title, subtitle, body if body else None, page_number=i+1)

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Create PPTX matching 集团公司 template style"
    )
    parser.add_argument("--output", "-o", required=True,
                        help="Output .pptx file path")
    parser.add_argument("--title", help="Presentation title (first slide)")
    parser.add_argument("--subtitle", help="Presentation subtitle")
    parser.add_argument("--slides", help="JSON file with slide definitions")
    parser.add_argument("--markdown", help="Markdown string or file path with slide content")
    parser.add_argument("--template", help="Path to a .pptx template (optional)")
    parser.add_argument("--end-slide", action="store_true", default=True,
                        help="Add end slide (default: True)")

    args = parser.parse_args()

    slides = []

    # Parse input
    if args.slides:
        with open(args.slides, "r", encoding="utf-8") as f:
            slides = json.load(f)
    elif args.markdown:
        if Path(args.markdown).exists():
            with open(args.markdown, "r", encoding="utf-8") as f:
                md_text = f.read()
        else:
            md_text = args.markdown
        slides = parse_markdown(md_text)
    else:
        # Single slide from --title
        slides = [{"layout": "title", "title": args.title or "汇报材料",
                   "subtitle": args.subtitle or ""}]

    if not slides:
        print("Error: no slides defined", file=sys.stderr)
        sys.exit(1)

    # Build
    output = create_presentation(slides, args.output, args.template)

    # Add end slide
    if args.end_slide:
        prs = Presentation(output)
        make_end_slide(prs)
        prs.save(output)

    print(f"PPTX saved: {output}")
    print(f"Slides: {len(slides) + (1 if args.end_slide else 0)}")


if __name__ == "__main__":
    main()
